"""Generate the THCC presentation deck as a .pptx for import into Google Slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------- design tokens ----------

BG       = RGBColor(0xFA, 0xF8, 0xF3)   # warm off-white
INK      = RGBColor(0x1A, 0x1A, 0x1A)   # near-black
ACCENT   = RGBColor(0xC4, 0x42, 0x2E)   # warm red
ACCENT2  = RGBColor(0x2E, 0x6E, 0xA0)   # cool blue
DIM      = RGBColor(0x6B, 0x6B, 0x6B)   # gray
LIGHT    = RGBColor(0xE6, 0xE2, 0xD8)   # tan-gray
CODE_BG  = RGBColor(0xF1, 0xEC, 0xE0)   # parchment
HIGHLIGHT= RGBColor(0xFF, 0xE8, 0x9C)   # pale yellow
DARK_BG  = RGBColor(0x14, 0x14, 0x14)   # near-black for the tease
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "Inter"
SERIF = "Georgia"
MONO = "Consolas"


# ---------- canvas ----------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------

def add_bg(slide, color=BG):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=SANS):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    tf.text = ""
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height,
             fill=None, line=INK, line_w=1.0, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_line(slide, x1, y1, x2, y2, color=INK, width=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def add_arrow(slide, x1, y1, x2, y2, color=INK, width=1.5):
    """Straight arrow from (x1,y1) to (x2,y2)."""
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    # add an arrowhead at end
    line_elem = ln.line._get_or_add_ln()
    tailEnd = etree.SubElement(line_elem, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    return ln


def add_code(slide, code, left, top, width, height,
             size=20, color=INK, bg=CODE_BG, highlights=None):
    """Add a monospace code block. highlights: list of line indices (0-based) to highlight."""
    if bg is not None:
        bgrect = add_rect(slide, left, top, width, height, fill=bg, line=None)
    tb = slide.shapes.add_textbox(
        left + Inches(0.18), top + Inches(0.12),
        width - Inches(0.36), height - Inches(0.24)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = MONO
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if highlights and i in highlights:
            # highlight the line by making it bold + colored
            run.font.bold = True
            run.font.color.rgb = ACCENT
    return tb


def slide_title(slide, text, color=DIM, size=14):
    """Small caps-ish title in the top-left corner."""
    add_text(slide, text.upper(), Inches(0.6), Inches(0.45),
             Inches(12.0), Inches(0.4),
             size=size, bold=True, color=color, font=SANS)


def slide_caption(slide, text, color=DIM, size=18, italic=True):
    """Italic caption line at the bottom."""
    add_text(slide, text, Inches(0.6), Inches(6.7),
             Inches(12.1), Inches(0.5),
             size=size, italic=italic, color=color,
             align=PP_ALIGN.CENTER, font=SANS)


def add_accumulator_box(slide, left, top, value, width=Inches(1.2), height=Inches(0.7),
                        label="acc"):
    add_rect(slide, left, top, width, height, fill=WHITE, line=INK, line_w=1.5)
    add_text(slide, value, left, top, width, height,
             size=22, bold=True, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
    add_text(slide, label, left, top + height + Inches(0.05),
             width, Inches(0.3), size=11, color=DIM,
             align=PP_ALIGN.CENTER, font=SANS)


# ============================================================
# SLIDE 1 — Opener
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_text(s, "2 + 2 = ?", Inches(0), Inches(2.4), SW, Inches(2.5),
         size=180, bold=True, color=INK, align=PP_ALIGN.CENTER, font=SERIF)
add_text(s, "ALU", Inches(0), Inches(5.3), SW, Inches(0.6),
         size=28, color=ACCENT, align=PP_ALIGN.CENTER, font=SANS, bold=True)
add_text(s, "(arithmetic logic unit)", Inches(0), Inches(5.85),
         SW, Inches(0.4),
         size=16, color=DIM, italic=True, align=PP_ALIGN.CENTER, font=SANS)


# ============================================================
# SLIDE 2 — Calculator + memory
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "ALU + Memory")

add_text(s, "(5 + 3)  ×  (7 + 2)",
         Inches(0), Inches(1.2), SW, Inches(1.0),
         size=56, bold=True, color=INK, align=PP_ALIGN.CENTER, font=SERIF)

# calculator block
calc_l, calc_t = Inches(2.8), Inches(3.0)
calc_w, calc_h = Inches(3.0), Inches(2.8)
add_rect(s, calc_l, calc_t, calc_w, calc_h, fill=LIGHT, line=INK, line_w=1.5, rounded=True)
add_text(s, "calculator", calc_l, calc_t + Inches(0.2), calc_w, Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)
# M+ / MR buttons
add_rect(s, calc_l + Inches(0.3), calc_t + Inches(0.9), Inches(1.0), Inches(0.6),
         fill=ACCENT, line=None, rounded=True)
add_text(s, "M+", calc_l + Inches(0.3), calc_t + Inches(0.9), Inches(1.0), Inches(0.6),
         size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)
add_rect(s, calc_l + Inches(1.7), calc_t + Inches(0.9), Inches(1.0), Inches(0.6),
         fill=ACCENT, line=None, rounded=True)
add_text(s, "MR", calc_l + Inches(1.7), calc_t + Inches(0.9), Inches(1.0), Inches(0.6),
         size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)
# screen showing current calc
add_rect(s, calc_l + Inches(0.3), calc_t + Inches(1.8), Inches(2.4), Inches(0.7),
         fill=WHITE, line=INK)
add_text(s, "8", calc_l + Inches(0.3), calc_t + Inches(1.8), Inches(2.4), Inches(0.7),
         size=28, bold=True, color=INK, font=MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# arrow to memory
add_arrow(s, calc_l + calc_w + Inches(0.1), calc_t + Inches(1.4),
          Inches(8.4), calc_t + Inches(1.4), color=DIM, width=2.0)
add_text(s, "stash", calc_l + calc_w + Inches(0.1), calc_t + Inches(0.9),
         Inches(2.0), Inches(0.4),
         size=14, italic=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)

# memory block
mem_l = Inches(8.4)
mem_t = calc_t + Inches(0.8)
add_rect(s, mem_l, mem_t, Inches(2.2), Inches(1.4), fill=WHITE, line=INK, line_w=1.5)
add_text(s, "memory", mem_l, mem_t + Inches(0.1), Inches(2.2), Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)
add_text(s, "8", mem_l, mem_t + Inches(0.5), Inches(2.2), Inches(0.8),
         size=44, bold=True, color=INK, font=MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

slide_caption(s, "stash an intermediate, recall it later — already half a CPU")


# ============================================================
# SLIDE 3 — The four pieces of a CPU
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "What makes a CPU")

# Four boxes laid out with labels
box_w, box_h = Inches(2.4), Inches(1.4)
y_top = Inches(2.0)
y_bot = Inches(4.4)

boxes = [
    ("Program\nmemory",     Inches(0.9),  y_top),
    ("Program\ncounter",    Inches(3.7),  y_top),
    ("Decoder",             Inches(6.5),  y_top),
    ("Clock",               Inches(9.3),  y_top),
    ("ALU",                 Inches(3.7),  y_bot),
    ("RAM",                 Inches(6.5),  y_bot),
]

for label, l, t in boxes:
    add_rect(s, l, t, box_w, box_h, fill=WHITE, line=INK, line_w=1.5, rounded=True)
    add_text(s, label, l, t, box_w, box_h,
             size=22, bold=True, color=INK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)

# Arrows showing flow
# pgm mem → decoder via PC
add_arrow(s, Inches(3.3),  y_top + Inches(0.7), Inches(3.7),  y_top + Inches(0.7), color=DIM, width=2)
add_arrow(s, Inches(6.1),  y_top + Inches(0.7), Inches(6.5),  y_top + Inches(0.7), color=DIM, width=2)
# decoder → ALU
add_arrow(s, Inches(7.7),  y_top + Inches(1.4), Inches(5.5),  y_bot, color=DIM, width=2)
# ALU ↔ RAM
add_arrow(s, Inches(6.1),  y_bot + Inches(0.7), Inches(6.5),  y_bot + Inches(0.7), color=DIM, width=2)
# clock to PC
add_arrow(s, Inches(9.3),  y_top + Inches(0.7), Inches(6.1),  y_top + Inches(0.7), color=ACCENT, width=2)

# footer ribbon
add_rect(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
         fill=ACCENT, line=None, rounded=True)
add_text(s, "fetch  →  decode  →  execute  →  advance",
         Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
         size=24, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)


# ============================================================
# SLIDE 4 — von Neumann architecture
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "The Von Neumann architecture")

# big memory rectangle with two regions
mem_l, mem_t = Inches(1.5), Inches(2.0)
mem_w, mem_h = Inches(6.5), Inches(3.5)
add_rect(s, mem_l, mem_t, mem_w, mem_h, fill=WHITE, line=INK, line_w=2.0, rounded=True)
add_text(s, "MEMORY", mem_l, mem_t + Inches(0.15), mem_w, Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)

# Program region
prog_l, prog_t = mem_l + Inches(0.3), mem_t + Inches(0.8)
prog_w = mem_w - Inches(0.6)
add_rect(s, prog_l, prog_t, prog_w, Inches(1.1), fill=HIGHLIGHT, line=DIM)
add_text(s, "program", prog_l, prog_t, prog_w, Inches(1.1),
         size=28, bold=True, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)

# Data region
data_l, data_t = prog_l, prog_t + Inches(1.3)
add_rect(s, data_l, data_t, prog_w, Inches(1.1), fill=LIGHT, line=DIM)
add_text(s, "data", data_l, data_t, prog_w, Inches(1.1),
         size=28, bold=True, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)

# CPU box on the right
cpu_l, cpu_t = Inches(9.5), Inches(3.2)
cpu_w, cpu_h = Inches(2.8), Inches(1.4)
add_rect(s, cpu_l, cpu_t, cpu_w, cpu_h, fill=ACCENT, line=None, rounded=True)
add_text(s, "CPU", cpu_l, cpu_t, cpu_w, cpu_h,
         size=36, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)

# arrow from program region into CPU
add_arrow(s, prog_l + prog_w, prog_t + Inches(0.55),
          cpu_l, cpu_t + Inches(0.4), color=ACCENT, width=2.5)
add_text(s, "reads its own instructions",
         Inches(7.8), Inches(2.7), Inches(2.2), Inches(0.4),
         size=12, italic=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)

# arrow from data region into CPU
add_arrow(s, data_l + prog_w, data_t + Inches(0.55),
          cpu_l, cpu_t + Inches(1.0), color=DIM, width=2)

slide_caption(s, "one memory.  programs and data both live there.")


# ============================================================
# SLIDE 5 — The full stack (preview)
# ============================================================
def draw_stack(slide, lit_layers=None):
    """Five-layer stack; lit_layers is a set of indices (0-4) to color."""
    if lit_layers is None:
        lit_layers = set()
    layers = [
        ("caesar.thcc",        "source",         ACCENT),
        ("AST",                "tree",           ACCENT2),
        (".asm  /  .hex",      "machine code",   ACCENT),
        ("THMM",               "hardware",       ACCENT2),
        ("RAM cells",          "result",         ACCENT),
    ]
    arrows = ["parse", "codegen + link", "load", "run", "read"]
    layer_w = Inches(7.0)
    layer_h = Inches(0.85)
    layer_l = Inches(3.2)
    top = Inches(1.4)
    gap = Inches(0.2)

    for i, (name, sub, col) in enumerate(layers):
        y = top + i * (layer_h + gap)
        is_lit = i in lit_layers
        fill = col if is_lit else WHITE
        text_color = WHITE if is_lit else INK
        sub_color = WHITE if is_lit else DIM
        add_rect(slide, layer_l, y, layer_w, layer_h,
                 fill=fill, line=INK if not is_lit else None,
                 line_w=1.5, rounded=True)
        add_text(slide, name, layer_l + Inches(0.4), y, Inches(3.5), layer_h,
                 size=22, bold=True, color=text_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        add_text(slide, sub, layer_l + Inches(3.5), y, Inches(3.0), layer_h,
                 size=14, italic=True, color=sub_color,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=SANS)
        # arrow + label
        if i < 4:
            ay = y + layer_h
            add_text(slide, arrows[i],
                     layer_l + layer_w + Inches(0.2), ay - Inches(0.05),
                     Inches(2.4), Inches(0.3),
                     size=12, italic=True, color=DIM, font=SANS)
            add_arrow(slide,
                      layer_l + layer_w / 2, ay,
                      layer_l + layer_w / 2, ay + gap,
                      color=DIM, width=2)


s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Today: every layer")
draw_stack(s, lit_layers=set())
slide_caption(s, "human-readable code  ↓  electrons (well, simulated)")


# ============================================================
# SLIDE 6 — Caesar tease
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_BG)
add_text(s, "BY THE END,  THIS CPU WILL READ THIS",
         Inches(0), Inches(0.6), SW, Inches(0.5),
         size=18, bold=True, color=DIM,
         align=PP_ALIGN.CENTER, font=SANS)
add_text(s, "JCJUM JS DCAFE",
         Inches(0), Inches(2.8), SW, Inches(2.0),
         size=120, bold=True, color=HIGHLIGHT,
         align=PP_ALIGN.CENTER, font=MONO)
add_text(s, "🔒",
         Inches(0), Inches(5.5), SW, Inches(1.0),
         size=72, color=ACCENT, align=PP_ALIGN.CENTER, font=SANS)


# ============================================================
# SLIDE 7 — Meet THMM
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Meet THMM")

# left: cpu schematic
cpu_l, cpu_t = Inches(0.8), Inches(1.5)
cpu_w, cpu_h = Inches(5.2), Inches(4.4)
add_rect(s, cpu_l, cpu_t, cpu_w, cpu_h, fill=WHITE, line=INK, line_w=2.0, rounded=True)
add_text(s, "THMM", cpu_l, cpu_t + Inches(0.15), cpu_w, Inches(0.5),
         size=20, bold=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)
# accumulator
add_rect(s, cpu_l + Inches(0.5), cpu_t + Inches(1.0), Inches(2.0), Inches(1.0),
         fill=ACCENT, line=None, rounded=True)
add_text(s, "ACC", cpu_l + Inches(0.5), cpu_t + Inches(1.0), Inches(2.0), Inches(1.0),
         size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)
# alu
add_rect(s, cpu_l + Inches(2.7), cpu_t + Inches(1.0), Inches(2.0), Inches(1.0),
         fill=LIGHT, line=INK, rounded=True)
add_text(s, "ALU", cpu_l + Inches(2.7), cpu_t + Inches(1.0), Inches(2.0), Inches(1.0),
         size=22, bold=True, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)
# ram strip
add_rect(s, cpu_l + Inches(0.5), cpu_t + Inches(2.5), Inches(4.2), Inches(1.5),
         fill=WHITE, line=INK, rounded=False)
add_text(s, "RAM[256]", cpu_l + Inches(0.5), cpu_t + Inches(2.5),
         Inches(4.2), Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)
# tiny grid of cells inside ram
for col in range(16):
    cx = cpu_l + Inches(0.6) + col * Inches(0.255)
    cy = cpu_t + Inches(3.0)
    add_rect(s, cx, cy, Inches(0.23), Inches(0.85), fill=LIGHT, line=DIM)

# right: instruction list
inst_l = Inches(7.0)
inst_t = Inches(1.5)
add_text(s, "9 instructions", inst_l, inst_t, Inches(5.5), Inches(0.5),
         size=18, bold=True, color=DIM, font=SANS)
inst_lines = [
    "loadm  addr     load mem → acc",
    "loadn  imm      load literal → acc",
    "store  addr     acc → mem",
    "addm   addr     acc += mem",
    "addn   imm      acc += literal",
    "subm   addr     acc -= mem",
    "mulm   addr     acc *= mem",
    "divm   addr     acc /= mem",
    "halt            stop",
]
for i, line in enumerate(inst_lines):
    add_text(s, line, inst_l, inst_t + Inches(0.6) + i * Inches(0.42),
             Inches(5.7), Inches(0.4),
             size=18, color=INK, font=MONO)

# footer ribbon
add_rect(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
         fill=ACCENT, line=None, rounded=True)
add_text(s, "16-bit  ·  256 cells  ·  9 instructions  ·  ~200 lines of Python",
         Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
         size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=SANS)


# ============================================================
# SLIDE 8 — The easy case
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "The easy case")

add_text(s, "c  =  a + b",
         Inches(0), Inches(1.0), SW, Inches(1.0),
         size=72, bold=True, color=INK, align=PP_ALIGN.CENTER, font=SERIF)

# big arrow down
add_arrow(s, SW / 2, Inches(2.2), SW / 2, Inches(2.7), color=DIM, width=3)

# code block (left-ish)
code = "loadm  a\naddm   b\nstore  c"
add_code(s, code, Inches(2.5), Inches(3.0), Inches(4.5), Inches(2.3), size=32)

# accumulator state on right
acc_l = Inches(8.0)
acc_t = Inches(3.2)
add_text(s, "accumulator", acc_l, acc_t - Inches(0.5),
         Inches(3.8), Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.CENTER, font=SANS)
add_accumulator_box(s, acc_l, acc_t, "?", label="start")
add_arrow(s, acc_l + Inches(1.3), acc_t + Inches(0.35),
          acc_l + Inches(1.6), acc_t + Inches(0.35), color=DIM, width=2)
add_accumulator_box(s, acc_l + Inches(1.6), acc_t, "a", label="loadm a")
add_arrow(s, acc_l + Inches(2.9), acc_t + Inches(0.35),
          acc_l + Inches(3.2), acc_t + Inches(0.35), color=DIM, width=2)
add_accumulator_box(s, acc_l + Inches(3.2), acc_t, "a+b", label="addm b")

slide_caption(s, "every expression's code leaves its value in the accumulator")


# ============================================================
# SLIDE 9 — The pain
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "The pain")

add_text(s, "(a + b)  ×  (c + d)",
         Inches(0), Inches(1.0), SW, Inches(1.0),
         size=64, bold=True, color=INK, align=PP_ALIGN.CENTER, font=SERIF)

code = (
    "loadm  a\n"
    "addm   b\n"
    "store  t0     ← stash\n"
    "loadm  c\n"
    "addm   d\n"
    "store  t1     ← stash\n"
    "loadm  t0     ← reload\n"
    "mulm   t1"
)
add_code(s, code, Inches(2.0), Inches(2.3), Inches(5.3), Inches(3.8),
         size=22, highlights=[2, 5, 6])

# right column: visualization
vis_l = Inches(8.2)
vis_t = Inches(2.5)
add_text(s, "memory", vis_l, vis_t - Inches(0.4), Inches(4.0), Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.LEFT, font=SANS)

cells = [("acc", "a+b → ... → answer", LIGHT),
         ("t0",  "a + b",               HIGHLIGHT),
         ("t1",  "c + d",               HIGHLIGHT)]
for i, (lab, val, col) in enumerate(cells):
    y = vis_t + i * Inches(1.0)
    add_rect(s, vis_l, y, Inches(1.5), Inches(0.7), fill=col, line=INK)
    add_text(s, lab, vis_l, y, Inches(1.5), Inches(0.7),
             size=20, bold=True, color=INK, font=MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, val, vis_l + Inches(1.7), y, Inches(3.0), Inches(0.7),
             size=18, color=INK, font=MONO,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

slide_caption(s, "now imagine 80 of these   —   by hand   —   without bugs",
              color=ACCENT, italic=True)


# ============================================================
# SLIDE 10 — Pythagoras source
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "A real program")

code = (
    "// 3-4-5 right triangle. Expected: hyp_sq = 25.\n"
    "\n"
    "int a = 4;\n"
    "int b = 3;\n"
    "int hyp_sq = a * a + b * b;"
)
add_code(s, code, Inches(2.5), Inches(2.0), Inches(8.3), Inches(3.5), size=32)

# small triangle in corner (upper right)
tri_l, tri_t = Inches(10.5), Inches(1.4)
# right triangle: legs 3 and 4 scaled
scale = Inches(0.5)
p1 = (tri_l, tri_t + 3 * scale)            # bottom-left
p2 = (tri_l + 4 * scale, tri_t + 3 * scale) # bottom-right
p3 = (tri_l, tri_t)                          # top-left
add_line(s, p1[0], p1[1], p2[0], p2[1], color=INK, width=2)
add_line(s, p2[0], p2[1], p3[0], p3[1], color=INK, width=2)
add_line(s, p3[0], p3[1], p1[0], p1[1], color=INK, width=2)
add_text(s, "3", tri_l - Inches(0.3), tri_t + Inches(0.6),
         Inches(0.3), Inches(0.4), size=14, color=DIM, font=SANS)
add_text(s, "4", tri_l + Inches(0.7), tri_t + Inches(1.6),
         Inches(0.5), Inches(0.4), size=14, color=DIM, font=SANS)
add_text(s, "?", tri_l + Inches(1.2), tri_t + Inches(0.5),
         Inches(0.5), Inches(0.4), size=18, bold=True,
         color=ACCENT, font=SERIF)

slide_caption(s, "three lines.  compile, run, read RAM[16].")


# ============================================================
# SLIDE 11 — The three squares (joke)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "What we got")

# Right triangle + three squares geometry
# Triangle vertices in slide coordinates:
ox, oy = Inches(5.6), Inches(4.5)   # right-angle corner (origin)
unit = Inches(0.45)                  # one unit of length
# Legs: a=4 along x (right), b=3 along y (up)
A = (ox, oy)                         # right-angle corner
B = (ox + 4 * unit, oy)              # along the 4-leg
C = (ox, oy - 3 * unit)              # along the 3-leg

# triangle edges
add_line(s, A[0], A[1], B[0], B[1], color=INK, width=3)
add_line(s, A[0], A[1], C[0], C[1], color=INK, width=3)
add_line(s, B[0], B[1], C[0], C[1], color=ACCENT, width=4)

# square on the 4-leg (below)
add_rect(s, A[0], A[1], 4 * unit, 4 * unit, fill=LIGHT, line=INK)
add_text(s, "16", A[0], A[1], 4 * unit, 4 * unit,
         size=44, bold=True, color=DIM, font=SERIF,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# square on the 3-leg (to the left)
add_rect(s, C[0] - 3 * unit, C[1], 3 * unit, 3 * unit, fill=LIGHT, line=INK)
add_text(s, "9", C[0] - 3 * unit, C[1], 3 * unit, 3 * unit,
         size=36, bold=True, color=DIM, font=SERIF,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# square on the hypotenuse — we approximate by drawing a tilted square
# easier: just label the hypotenuse area as a separate big square off to the right
hyp_l = Inches(8.6)
hyp_t = Inches(1.6)
hyp_s = Inches(2.5)
add_rect(s, hyp_l, hyp_t, hyp_s, hyp_s, fill=HIGHLIGHT, line=ACCENT, line_w=3)
add_text(s, "25", hyp_l, hyp_t, hyp_s, hyp_s,
         size=84, bold=True, color=ACCENT, font=SERIF,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, "the square on the hypotenuse",
         hyp_l - Inches(0.3), hyp_t + hyp_s + Inches(0.1),
         hyp_s + Inches(0.6), Inches(0.4),
         size=14, italic=True, color=DIM, font=SANS, align=PP_ALIGN.CENTER)
# arrow from hypotenuse to the big square
add_arrow(s, B[0], (B[1] + C[1]) / 2, hyp_l, hyp_t + hyp_s / 2,
          color=ACCENT, width=2)

# the "no sqrt" callout
nosqrt_l = Inches(0.8)
nosqrt_t = Inches(2.0)
add_rect(s, nosqrt_l, nosqrt_t, Inches(3.2), Inches(1.8),
         fill=WHITE, line=ACCENT, line_w=3, rounded=True)
add_text(s, "sqrt", nosqrt_l, nosqrt_t + Inches(0.4),
         Inches(3.2), Inches(0.7),
         size=44, bold=True, color=DIM, font=MONO,
         align=PP_ALIGN.CENTER)
# big red X — two diagonal lines across the box
add_line(s, nosqrt_l + Inches(0.3), nosqrt_t + Inches(0.3),
         nosqrt_l + Inches(2.9), nosqrt_t + Inches(1.5),
         color=ACCENT, width=8)
add_line(s, nosqrt_l + Inches(2.9), nosqrt_t + Inches(0.3),
         nosqrt_l + Inches(0.3), nosqrt_t + Inches(1.5),
         color=ACCENT, width=8)
add_text(s, "not in the instruction set",
         nosqrt_l, nosqrt_t + Inches(1.9),
         Inches(3.2), Inches(0.4),
         size=14, italic=True, color=DIM, font=SANS, align=PP_ALIGN.CENTER)

slide_caption(s, "we can compute the area.   the side is not our problem.")


# ============================================================
# SLIDE 12 — Parse → tree
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Step 1   ·   Parse  →  tree")

# source line
add_text(s, "int hyp_sq  =  a * a + b * b ;",
         Inches(0), Inches(1.3), SW, Inches(0.8),
         size=36, bold=True, color=INK, font=MONO,
         align=PP_ALIGN.CENTER)

# arrow down
add_arrow(s, SW / 2, Inches(2.3), SW / 2, Inches(2.8), color=DIM, width=3)

# tree
def tree_node(slide, label, cx, cy, w=Inches(1.0), h=Inches(0.7), fill=WHITE):
    add_rect(slide, cx - w / 2, cy, w, h, fill=fill, line=INK, line_w=1.5, rounded=True)
    add_text(slide, label, cx - w / 2, cy, w, h,
             size=24, bold=True, color=INK, font=MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return (cx, cy, cx, cy + h)  # top-center, bottom-center

# tree positions
root_x = SW / 2
root_y = Inches(3.0)
mul_y = Inches(4.3)
leaf_y = Inches(5.6)

# spread
half = Inches(2.6)

n_root = tree_node(s, "+", root_x, root_y, fill=HIGHLIGHT)
n_mulL = tree_node(s, "×", root_x - half, mul_y)
n_mulR = tree_node(s, "×", root_x + half, mul_y)

leaf_off = Inches(1.0)
n_aL = tree_node(s, "a", root_x - half - leaf_off, leaf_y, fill=LIGHT)
n_aR = tree_node(s, "a", root_x - half + leaf_off, leaf_y, fill=LIGHT)
n_bL = tree_node(s, "b", root_x + half - leaf_off, leaf_y, fill=LIGHT)
n_bR = tree_node(s, "b", root_x + half + leaf_off, leaf_y, fill=LIGHT)

# edges
def tree_edge(s, parent_node, child_node):
    px, py, pbx, pby = parent_node
    cx, cy, cbx, cby = child_node
    add_line(s, pbx, pby, cx, cy, color=DIM, width=1.5)

tree_edge(s, n_root, n_mulL)
tree_edge(s, n_root, n_mulR)
tree_edge(s, n_mulL, n_aL)
tree_edge(s, n_mulL, n_aR)
tree_edge(s, n_mulR, n_bL)
tree_edge(s, n_mulR, n_bR)

slide_caption(s, "expressions nest, so they're trees")


# ============================================================
# SLIDE 13 — H1: AST in Haskell
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Haskell  ·  the abstract syntax tree")

code = (
    "data Expr\n"
    "    = Lit   Int\n"
    "    | Var   String\n"
    "    | BinOp Op Expr Expr\n"
    "\n"
    "data Op = Add | Sub | Mul | Div"
)
add_code(s, code, Inches(2.5), Inches(2.0), Inches(8.3), Inches(3.5),
         size=28, highlights=[3])

slide_caption(s, "three lines define the entire expression language")


# ============================================================
# SLIDE 14 — H2: Operator table
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Haskell  ·  precedence is data")

code = (
    "operatorTable =\n"
    "    [ [ binaryL \"*\" Mul, binaryL \"/\" Div ]   ← binds tighter\n"
    "    , [ binaryL \"+\" Add, binaryL \"-\" Sub ]   ← binds looser\n"
    "    ]"
)
add_code(s, code, Inches(1.5), Inches(2.4), Inches(10.3), Inches(2.6), size=24,
         highlights=[1, 2])

slide_caption(s, "row order  =  precedence.   grammar as data, not code.")


# ============================================================
# SLIDE 15 — Codegen: easy
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Codegen  ·  easy case")

# tree fragment for a*a
tx = Inches(2.8)
ty = Inches(1.8)
n_mul = tree_node(s, "×", tx, ty)
n_a1 = tree_node(s, "a", tx - Inches(0.9), ty + Inches(1.3), fill=LIGHT)
n_a2 = tree_node(s, "a", tx + Inches(0.9), ty + Inches(1.3), fill=LIGHT)
tree_edge(s, n_mul, n_a1)
tree_edge(s, n_mul, n_a2)

# arrow
add_arrow(s, Inches(4.5), Inches(2.6), Inches(5.4), Inches(2.6),
          color=DIM, width=3)

# code
code = "loadm  a\nmulm   a"
add_code(s, code, Inches(5.6), Inches(2.0), Inches(3.3), Inches(1.5), size=32)

# accumulator on right
acc_l = Inches(9.5)
acc_t = Inches(2.0)
add_text(s, "accumulator", acc_l, acc_t - Inches(0.4),
         Inches(3.3), Inches(0.4),
         size=14, bold=True, color=DIM, align=PP_ALIGN.LEFT, font=SANS)
add_accumulator_box(s, acc_l, acc_t, "?", label="start")
add_accumulator_box(s, acc_l, acc_t + Inches(1.2), "a", label="loadm a")
add_accumulator_box(s, acc_l, acc_t + Inches(2.4), "a²", label="mulm a")

slide_caption(s, "leaves emit one instruction.   the contract: value lives in the accumulator.")


# ============================================================
# SLIDE 16 — Codegen: temp dance
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Codegen  ·  the temp dance")

# tree at top
tx = Inches(3.0)
ty = Inches(1.3)
n_root = tree_node(s, "+", tx, ty, fill=HIGHLIGHT)
n_mulL = tree_node(s, "×", tx - Inches(1.4), ty + Inches(1.1))
n_mulR = tree_node(s, "×", tx + Inches(1.4), ty + Inches(1.1))
n_a1 = tree_node(s, "a", tx - Inches(1.9), ty + Inches(2.2), fill=LIGHT)
n_a2 = tree_node(s, "a", tx - Inches(0.9), ty + Inches(2.2), fill=LIGHT)
n_b1 = tree_node(s, "b", tx + Inches(0.9), ty + Inches(2.2), fill=LIGHT)
n_b2 = tree_node(s, "b", tx + Inches(1.9), ty + Inches(2.2), fill=LIGHT)
tree_edge(s, n_root, n_mulL)
tree_edge(s, n_root, n_mulR)
tree_edge(s, n_mulL, n_a1)
tree_edge(s, n_mulL, n_a2)
tree_edge(s, n_mulR, n_b1)
tree_edge(s, n_mulR, n_b2)

# code
code = (
    "loadm  a\n"
    "mulm   a\n"
    "store  t0     ← stash a²\n"
    "loadm  b\n"
    "mulm   b\n"
    "store  t1     ← stash b²\n"
    "loadm  t0     ← reload\n"
    "addm   t1\n"
    "store  hyp_sq"
)
add_code(s, code, Inches(6.0), Inches(1.6), Inches(4.5), Inches(4.6),
         size=20, highlights=[2, 5, 6])

# right column: temps
acc_l = Inches(11.0)
acc_t = Inches(2.0)
add_rect(s, acc_l, acc_t, Inches(1.5), Inches(0.7), fill=ACCENT, line=None, rounded=True)
add_text(s, "acc", acc_l, acc_t, Inches(1.5), Inches(0.7),
         size=18, bold=True, color=WHITE, font=MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, acc_l, acc_t + Inches(1.0), Inches(1.5), Inches(0.7), fill=HIGHLIGHT, line=INK)
add_text(s, "t0", acc_l, acc_t + Inches(1.0), Inches(1.5), Inches(0.7),
         size=18, bold=True, color=INK, font=MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_rect(s, acc_l, acc_t + Inches(2.0), Inches(1.5), Inches(0.7), fill=HIGHLIGHT, line=INK)
add_text(s, "t1", acc_l, acc_t + Inches(2.0), Inches(1.5), Inches(0.7),
         size=18, bold=True, color=INK, font=MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

slide_caption(s, "stash the left, compute the right, recombine")


# ============================================================
# SLIDE 17 — H3: genExpr
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Haskell  ·  the codegen function")

code = (
    "genExpr _    (Lit n)        = ...   -- emit loadn\n"
    "genExpr vars (Var x)        = ...   -- emit loadm\n"
    "genExpr vars (BinOp op l r) = ...   -- recurse: l, then r, combine"
)
add_code(s, code, Inches(0.6), Inches(2.6), Inches(12.1), Inches(2.4), size=22)

slide_caption(s, "pattern-matches on tree shape.   leaves emit; branches recurse.")


# ============================================================
# SLIDE 18 — Link
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Step 3   ·   Link  ·  names → addresses")

# left column: before
left_code = (
    "loadm  a\n"
    "mulm   a\n"
    "store  t0\n"
    "loadm  b\n"
    "mulm   b\n"
    "store  t1\n"
    "loadm  t0\n"
    "addm   t1\n"
    "store  hyp_sq\n"
    "halt"
)
add_text(s, "before", Inches(1.0), Inches(1.4), Inches(4.0), Inches(0.4),
         size=16, bold=True, color=DIM, align=PP_ALIGN.LEFT, font=SANS)
add_code(s, left_code, Inches(1.0), Inches(1.8), Inches(4.0), Inches(4.4), size=18)

# arrow
add_arrow(s, Inches(5.4), Inches(4.0), Inches(6.0), Inches(4.0),
          color=ACCENT, width=3)

# right column: after
right_code = (
    " 4: loadm  14\n"
    " 5: mulm   14\n"
    " 6: store  17\n"
    " 7: loadm  15\n"
    " 8: mulm   15\n"
    " 9: store  18\n"
    "10: loadm  17\n"
    "11: addm   18\n"
    "12: store  16\n"
    "13: halt"
)
add_text(s, "after", Inches(6.2), Inches(1.4), Inches(4.0), Inches(0.4),
         size=16, bold=True, color=DIM, align=PP_ALIGN.LEFT, font=SANS)
add_code(s, right_code, Inches(6.2), Inches(1.8), Inches(4.0), Inches(4.4), size=18)

# address table on far right
tab_l = Inches(10.5)
tab_t = Inches(1.8)
add_text(s, "address table", tab_l, tab_t - Inches(0.4),
         Inches(2.5), Inches(0.4),
         size=14, bold=True, color=DIM, font=SANS)
table_rows = [
    ("a",      "14"),
    ("b",      "15"),
    ("hyp_sq", "16"),
    ("t0",     "17"),
    ("t1",     "18"),
]
add_rect(s, tab_l, tab_t, Inches(2.5), Inches(0.4 * len(table_rows) + 0.2),
         fill=LIGHT, line=DIM)
for i, (name, addr) in enumerate(table_rows):
    add_text(s, name, tab_l + Inches(0.2), tab_t + Inches(0.1) + i * Inches(0.4),
             Inches(1.2), Inches(0.4),
             size=18, color=INK, font=MONO)
    add_text(s, "→", tab_l + Inches(1.3), tab_t + Inches(0.1) + i * Inches(0.4),
             Inches(0.3), Inches(0.4),
             size=18, color=DIM, font=MONO)
    add_text(s, addr, tab_l + Inches(1.6), tab_t + Inches(0.1) + i * Inches(0.4),
             Inches(0.7), Inches(0.4),
             size=18, bold=True, color=ACCENT, font=MONO)


# ============================================================
# SLIDE 19 — Four panes
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Same program.   four ways.")

# divide screen 2x2
mid_x = SW / 2
top_y = Inches(1.2)
mid_y = Inches(4.2)
pane_w = SW / 2 - Inches(0.2)
pane_h_top = Inches(2.9)
pane_h_bot = Inches(3.0)

# Top-left: source
src_code = (
    "int a = 4;\n"
    "int b = 3;\n"
    "int hyp_sq = a * a + b * b;"
)
add_text(s, "source  (.thcc)", Inches(0.6), top_y - Inches(0.05),
         pane_w, Inches(0.3), size=12, bold=True, color=DIM, font=SANS)
add_code(s, src_code, Inches(0.6), top_y + Inches(0.25),
         pane_w, pane_h_top - Inches(0.3), size=18)

# Top-right: AST
add_text(s, "AST", mid_x + Inches(0.1), top_y - Inches(0.05),
         pane_w, Inches(0.3), size=12, bold=True, color=DIM, font=SANS)
ast_l = mid_x + Inches(0.1)
ast_t = top_y + Inches(0.4)
ast_w = pane_w
ast_cx = ast_l + ast_w / 2
n_root = tree_node(s, "+", ast_cx, ast_t, fill=HIGHLIGHT)
n_mulL = tree_node(s, "×", ast_cx - Inches(1.2), ast_t + Inches(0.95))
n_mulR = tree_node(s, "×", ast_cx + Inches(1.2), ast_t + Inches(0.95))
n_a1 = tree_node(s, "a", ast_cx - Inches(1.6), ast_t + Inches(1.9), fill=LIGHT)
n_a2 = tree_node(s, "a", ast_cx - Inches(0.8), ast_t + Inches(1.9), fill=LIGHT)
n_b1 = tree_node(s, "b", ast_cx + Inches(0.8), ast_t + Inches(1.9), fill=LIGHT)
n_b2 = tree_node(s, "b", ast_cx + Inches(1.6), ast_t + Inches(1.9), fill=LIGHT)
tree_edge(s, n_root, n_mulL)
tree_edge(s, n_root, n_mulR)
tree_edge(s, n_mulL, n_a1)
tree_edge(s, n_mulL, n_a2)
tree_edge(s, n_mulR, n_b1)
tree_edge(s, n_mulR, n_b2)

# Bottom-left: assembly
asm_code = (
    " 0: loadn  4\n"
    " 1: store  14\n"
    " 2: loadn  3\n"
    " 3: store  15\n"
    " 4: loadm  14\n"
    " 5: mulm   14\n"
    " 6: store  17\n"
    " 7: loadm  15\n"
    " 8: mulm   15\n"
    " 9: store  18\n"
    "10: loadm  17\n"
    "11: addm   18\n"
    "12: store  16\n"
    "13: halt"
)
add_text(s, "assembly  (.asm)", Inches(0.6), mid_y - Inches(0.05),
         pane_w, Inches(0.3), size=12, bold=True, color=DIM, font=SANS)
add_code(s, asm_code, Inches(0.6), mid_y + Inches(0.25),
         pane_w, pane_h_bot - Inches(0.3), size=13)

# Bottom-right: hex / bits
hex_code = (
    "0x3004    0011 0000 0000 0100\n"
    "0x400E    0100 0000 0000 1110\n"
    "0x3003    0011 0000 0000 0011\n"
    "0x400F    0100 0000 0000 1111\n"
    "0x200E    0010 0000 0000 1110\n"
    "0xB00E    1011 0000 0000 1110\n"
    "0x4011    0100 0000 0001 0001\n"
    "0x200F    0010 0000 0000 1111\n"
    "0xB00F    1011 0000 0000 1111\n"
    "0x4012    0100 0000 0001 0010\n"
    "0x2011    0010 0000 0001 0001\n"
    "0x7012    0111 0000 0001 0010\n"
    "0x4010    0100 0000 0001 0000\n"
    "0x1000    0001 0000 0000 0000"
)
add_text(s, "machine code  (.hex / .bits)",
         mid_x + Inches(0.1), mid_y - Inches(0.05),
         pane_w, Inches(0.3), size=12, bold=True, color=DIM, font=SANS)
add_code(s, hex_code, mid_x + Inches(0.1), mid_y + Inches(0.25),
         pane_w, pane_h_bot - Inches(0.3), size=11)


# ============================================================
# SLIDE 20 — Caesar fun facts
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Caesar Cipher")

bullets_l = Inches(0.8)
bullets_t = Inches(2.0)
bullet_lines = [
    ("Used in Roman military dispatches.",  "~50 BC"),
    ("Each letter shifted by a fixed key.", "shift 3:  A → D, B → E, C → F"),
    ("Trivially broken by frequency analysis.", "but it's the ancestor of every cipher since"),
    ("ROT13 (shift 13) is still used today.",  "spoiler tags on the internet"),
]
for i, (head, sub) in enumerate(bullet_lines):
    y = bullets_t + i * Inches(1.0)
    # bullet dot
    add_rect(s, bullets_l, y + Inches(0.18), Inches(0.18), Inches(0.18),
             fill=ACCENT, line=None, rounded=True)
    add_text(s, head, bullets_l + Inches(0.4), y, Inches(7.0), Inches(0.5),
             size=24, bold=True, color=INK, font=SANS)
    add_text(s, sub, bullets_l + Inches(0.4), y + Inches(0.45), Inches(7.0), Inches(0.4),
             size=16, italic=True, color=DIM, font=SANS)

# alphabet wheel on right
wheel_cx = Inches(10.5)
wheel_cy = Inches(4.2)
outer_r = Inches(1.6)
inner_r = Inches(1.0)
add_rect(s, wheel_cx - outer_r, wheel_cy - outer_r, 2 * outer_r, 2 * outer_r,
         fill=LIGHT, line=INK, line_w=2)
# we'll fake "wheel" with a rounded square; mark a few letters and shift arrow
add_text(s, "A", wheel_cx - outer_r, wheel_cy - outer_r, 2 * outer_r, Inches(0.4),
         size=18, bold=True, color=INK, font=SERIF, align=PP_ALIGN.CENTER)
add_text(s, "N", wheel_cx - outer_r, wheel_cy + outer_r - Inches(0.4),
         2 * outer_r, Inches(0.4),
         size=18, bold=True, color=INK, font=SERIF, align=PP_ALIGN.CENTER)
add_text(s, "G", wheel_cx - outer_r - Inches(0.05), wheel_cy - Inches(0.2),
         Inches(0.4), Inches(0.4),
         size=18, bold=True, color=INK, font=SERIF, align=PP_ALIGN.CENTER)
add_text(s, "T", wheel_cx + outer_r - Inches(0.4), wheel_cy - Inches(0.2),
         Inches(0.4), Inches(0.4),
         size=18, bold=True, color=INK, font=SERIF, align=PP_ALIGN.CENTER)
# inner ring
add_rect(s, wheel_cx - inner_r, wheel_cy - inner_r, 2 * inner_r, 2 * inner_r,
         fill=ACCENT, line=None)
add_text(s, "shift", wheel_cx - inner_r, wheel_cy - Inches(0.4),
         2 * inner_r, Inches(0.4),
         size=14, bold=True, color=WHITE, font=SANS, align=PP_ALIGN.CENTER)
add_text(s, "3", wheel_cx - inner_r, wheel_cy,
         2 * inner_r, Inches(0.7),
         size=44, bold=True, color=WHITE, font=SERIF, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 21 — The stack, all lit
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
slide_title(s, "Every time you press Run")
draw_stack(s, lit_layers={0, 1, 2, 3, 4})
slide_caption(s, "parse  ·  generate  ·  link  ·  execute")


# ---------- save ----------

out = "THCC_presentation.pptx"
prs.save(out)
print(f"wrote {out}: {len(prs.slides)} slides")
